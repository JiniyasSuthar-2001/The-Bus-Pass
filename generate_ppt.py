from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define some brand colors if possible (though limited without a template)
    # We'll just stick to standard styling but ensure text is clear.

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Bus Pass Barcode System"
    subtitle.text = "Modernizing Public Transport Ticketing\nPresented by: Developer"

    # Slide 2: Introduction / Problem & Solution
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Introduction"
    content = slide.placeholders[1]
    
    text_frame = content.text_frame
    text_frame.text = "Problem:"
    p = text_frame.add_paragraph()
    p.text = "• Traditional paper ticketing is inefficient and prone to loss."
    p.level = 1
    p = text_frame.add_paragraph()
    p.text = "• Long queues for buying passes and tickets."
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "\nSolution:"
    p.font.bold = True
    p = text_frame.add_paragraph()
    p.text = "• A digital platform for managing bus passes and tickets."
    p.level = 1
    p = text_frame.add_paragraph()
    p.text = "• QR/Barcode-based verification for conductors."
    p.level = 1
    p = text_frame.add_paragraph()
    p.text = "• Seamless wallet-based payments."
    p.level = 1

    # Slide 3: Key Features
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Key Features"
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    bullets = [
        "User Roles: Admin, User, Conductor.",
        "Digital Wallet: Easy recharge and payments.",
        "Smart Routing: Search routes by Source & Destination.",
        "Instant Booking: Buy tickets for specific bus schedules.",
        "Barcode Verification: Secure scanning by Conductors.",
        "Pass Management: Monthly/Weekly passes linked to user profile."
    ]
    
    for bullet in bullets:
        p = text_frame.add_paragraph()
        p.text = bullet
        p.space_after = Pt(10)

    # Slide 4: User Journey
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "User Journey"
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    steps = [
        "1. Registration: User creates an account.",
        "2. Wallet Recharge: Adds money to digital wallet.",
        "3. Route Search: Finds available buses for desired route.",
        "4. Purchase: Buys a ticket or pass using wallet balance.",
        "5. Boarding: Shows digital pass/ticket to Conductor.",
        "6. Verification: Conductor scans barcode to validate."
    ]
    
    for step in steps:
        p = text_frame.add_paragraph()
        p.text = step
        p.space_after = Pt(10)

    # Slide 5: Technical Stack
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Technical Stack"
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    text_frame.text = "Backend Framework:"
    p = text_frame.add_paragraph()
    p.text = "• Django (Python)"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "Database:"
    p = text_frame.add_paragraph()
    p.text = "• SQLite (Development) / PostgreSQL (Production)"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "Frontend:"
    p = text_frame.add_paragraph()
    p.text = "• HTML5, CSS3, JavaScript (Django Templates)"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "Libraries:"
    p = text_frame.add_paragraph()
    p.text = "• python-barcode (Generation), Razorpay (Payments)"
    p.level = 1

    # Slide 6: Future Enhancements
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Future Scope"
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    bullets = [
        "Live Bus Tracking using GPS Integration.",
        "Mobile App (Android/iOS) for better accessibility.",
        "Automated Fare Calculation based on distance.",
        "Loyalty points and discounts for frequent travelers."
    ]
    
    for bullet in bullets:
        p = text_frame.add_paragraph()
        p.text = bullet
        p.space_after = Pt(10)

    # Slide 7: Conclusion
    slide_layout = prs.slide_layouts[0] # Title Slide style for end
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Thank You"
    subtitle.text = "Questions & Feedback?"

    output_file = "Bus_Pass_System_Presentation.pptx"
    prs.save(output_file)
    print(f"Presentation saved to: {output_file}")

if __name__ == "__main__":
    create_presentation()
